#!/usr/bin/env python3
"""
Generate comprehensive PowerPoint presentation for IMS (Inventory Management System)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(31, 78, 121)  # Professional blue
ACCENT_COLOR = RGBColor(192, 0, 0)    # Red accent
TEXT_COLOR = RGBColor(51, 51, 51)     # Dark gray
LIGHT_BG = RGBColor(242, 242, 242)   # Light gray

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(192, 0, 0)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list=None, subheading=None):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title background bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_fill = title_shape.fill
    title_fill.solid()
    title_fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add subheading if provided
    if subheading:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
        sub_frame = sub_box.text_frame
        p = sub_frame.paragraphs[0]
        p.text = subheading
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = ACCENT_COLOR
    
    # Add content
    if content_list:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            
            # Add bullet point
            if item.startswith("• "):
                p.level = 0
            elif item.startswith("  "):
                p.level = 1
                p.font.size = Pt(16)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title background bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_fill = title_shape.fill
    title_fill.solid()
    title_fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(6))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(6))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, 
    "INVENTORY MANAGEMENT SYSTEM", 
    "Enterprise-Grade Stock Management Solution")

# Slide 2: System Overview
add_content_slide(prs, 
    "System Overview",
    [
        "✓ Enterprise-grade web-based inventory platform",
        "✓ Multi-location inventory management with hierarchical control",
        "✓ Integrated procurement and stock issuance workflows",
        "✓ Real-time inventory tracking and reconciliation",
        "✓ Multi-level approval workflows with role-based access",
        "✓ Production-ready system managing 15+ items across multiple locations"
    ])

# Slide 3: Technology Stack
add_two_column_slide(prs,
    "Technology Stack",
    [
        "Frontend:",
        "• React 18 + TypeScript",
        "• Vite (build tool)",
        "• shadcn/ui Components",
        "• Tailwind CSS",
        "",
        "Authentication:",
        "• AspNetCore Identity",
        "• SSO Support"
    ],
    [
        "Backend:",
        "• Node.js + Express",
        "• CommonJS Modules",
        "• Winston (Logging)",
        "• Multer (File Uploads)",
        "",
        "Database:",
        "• SQL Server 2022",
        "• 50+ Tables"
    ])

# Slide 4: Key Modules - Part 1
add_content_slide(prs,
    "Core Modules",
    [
        "1. STOCK ACQUISITION (Procurement)",
        "   • Tender management (Contract, Spot-Purchase, Annual)",
        "   • Vendor evaluation and bidding",
        "   • Purchase order creation and tracking",
        "",
        "2. STOCK ISSUANCE",
        "   • Item request system with custom quantities",
        "   • 3-level approval workflow",
        "   • Per-item decision making (Approve/Return/Reject)",
        "   • Automatic inventory deduction"
    ])

# Slide 5: Key Modules - Part 2
add_content_slide(prs,
    "Core Modules (Continued)",
    [
        "3. INVENTORY VERIFICATION",
        "   • Physical stock verification by store keepers",
        "   • Verification request assignment",
        "   • Reconciliation with system records",
        "",
        "4. ANNUAL TENDER FRAMEWORK",
        "   • Year-long vendor contracts",
        "   • Item group organization",
        "   • Per-item vendor assignment",
        "",
        "5. HIERARCHICAL INVENTORY",
        "   • Admin Store (Central warehouse)",
        "   • Wing Store (Department-level)",
        "   • Personal Store (User-allocated items)"
    ])

# Slide 6: Key Workflows
add_content_slide(prs,
    "Main Operational Workflows",
    [
        "Stock Acquisition:",
        "Request → Tender → Vendor Bids → Award → Delivery → Distribution",
        "",
        "Stock Issuance:",
        "Request → Supervisor Approval → Finance → Admin → Issue → Deduct",
        "",
        "Inventory Verification:",
        "Request → Store Keeper → Physical Count → Reconciliation → Update",
        "",
        "Annual Tender:",
        "Define Groups → Create Framework → Assign Vendors → POs → Deliveries"
    ])

# Slide 7: User Roles & Permissions
add_content_slide(prs,
    "User Roles & Permissions",
    [
        "Super Admin - Full system access (50 permissions)",
        "Admin - Inventory, Users, Settings management",
        "Wing Supervisor - Wing inventory and approvals",
        "Store Keeper - Physical inventory verification",
        "Finance/Approver - Purchase order and financial approvals",
        "Users/Requesters - Item requests and personal allocations",
        "",
        "Total Active Users: 499+ registered across the system"
    ])

# Slide 8: Database Architecture
add_two_column_slide(prs,
    "Database Architecture",
    [
        "Master Data Tables:",
        "• Item Masters (15+ items)",
        "• Categories & Subcategories",
        "• Vendors (7+ vendors)",
        "",
        "Procurement Tables:",
        "• Tenders & Tender Items",
        "• Vendors & Bidding",
        "• Purchase Orders",
        "• Deliveries & Serial Numbers"
    ],
    [
        "Inventory Tables:",
        "• Stock Levels (3-tier)",
        "• Issuance Requests",
        "• Stock Returns",
        "• Current Inventory Stock",
        "",
        "Control Tables:",
        "• Approvals & Workflows",
        "• Approval History",
        "• User Designations",
        "• Organizational Hierarchy"
    ])

# Slide 9: System Features
add_content_slide(prs,
    "Key System Features",
    [
        "✅ Real-time Stock Tracking",
        "✅ Multi-Step Hierarchical Approval Workflows",
        "✅ Complete Audit Trail & History",
        "✅ Role-Based Access Control (RBAC)",
        "✅ Three-Level Inventory Management",
        "✅ Multiple Tender Types Support",
        "✅ Digital Annual Framework Contracts",
        "✅ Serial Number & Equipment Tracking",
        "✅ Document Upload & Attachment Support",
        "✅ Soft Delete with Data Recovery"
    ])

# Slide 10: Frontend Architecture
add_content_slide(prs,
    "Frontend Architecture",
    [
        "Core Pages & Dashboards:",
        "• Approval Dashboard (Supervisor/Finance)",
        "• Main Dashboard",
        "• Stock Issuance Management",
        "• Current Inventory Stock View",
        "• Tender Creation & Management",
        "• Store Keeper Verification",
        "• User Role Assignment",
        "• Purchase Order Dashboard",
        "",
        "Component Library: 25+ Reusable UI Components"
    ])

# Slide 11: Backend Architecture
add_content_slide(prs,
    "Backend Architecture",
    [
        "API Server (Express.js):",
        "• 40+ REST Endpoints",
        "• Authentication & Authorization Middleware",
        "• Comprehensive Logging (Winston)",
        "• Error Handling & Validation",
        "",
        "Key API Routes:",
        "• /api/tenders - Tender operations",
        "• /api/stock-issuance - Issuance management",
        "• /api/approvals - Approval workflows",
        "• /api/inventory - Stock queries",
        "• /api/purchase-orders - PO management",
        "• /api/users - User management"
    ])

# Slide 12: Project Status & Completion
add_two_column_slide(prs,
    "Project Status",
    [
        "✅ COMPLETED:",
        "• Core Inventory",
        "• Stock Issuance",
        "• Approval Workflows",
        "• Tender Management",
        "• Annual Tenders",
        "• Purchase Orders",
        "• Verification System",
        "• User/Role Management"
    ],
    [
        "🟡 IN PROGRESS:",
        "• Reporting & Analytics",
        "",
        "🔵 FUTURE ENHANCEMENTS:",
        "• Mobile Application",
        "• Advanced Analytics",
        "• Real-time Dashboards",
        "• AI-based Recommendations"
    ])

# Slide 13: Deployment & Configuration
add_content_slide(prs,
    "Deployment & Configuration",
    [
        "Development Stack:",
        "• Frontend Port: 8080 (Vite dev server)",
        "• Backend Port: 3001 (Express API)",
        "",
        "Configuration Management:",
        "• Environment-specific configs (dev/test/staging/prod)",
        "• Docker Support (Dockerfile + docker-compose)",
        "• Database Migrations (100+ SQL scripts)",
        "",
        "Deployment Automation:",
        "• PowerShell deployment scripts",
        "• Bash automation support"
    ])

# Slide 14: Data Security & Compliance
add_content_slide(prs,
    "Security & Data Management",
    [
        "Authentication & Authorization:",
        "✓ AspNetCore Identity with SSO",
        "✓ Role-Based Access Control (RBAC)",
        "✓ Multi-level approval workflows",
        "",
        "Data Management:",
        "✓ Soft delete mechanism (no permanent data loss)",
        "✓ Complete audit trail & approval history",
        "✓ Transaction-based operations",
        "✓ SQL Server security features"
    ])

# Slide 15: Business Impact
add_content_slide(prs,
    "Business Impact & Benefits",
    [
        "Operational Efficiency:",
        "→ Automated procurement-to-inventory workflow",
        "→ Real-time stock visibility across 3 inventory levels",
        "",
        "Control & Compliance:",
        "→ Multi-step approval ensuring accountability",
        "→ Complete audit trail for compliance",
        "",
        "Scalability:",
        "→ Supports 499+ users across organization",
        "→ Handles 15+ item types with unlimited expansion",
        "→ Multi-location inventory management"
    ])

# Slide 16: Implementation Highlights
add_content_slide(prs,
    "Implementation Highlights",
    [
        "✓ Successfully integrated 3 inventory levels",
        "✓ Implemented complex approval workflows",
        "✓ Annual tender system with framework contracts",
        "✓ Real-time stock tracking & reconciliation",
        "✓ Vendor management & bidding system",
        "✓ Purchase order automation",
        "✓ Serial number tracking for equipment",
        "✓ Comprehensive user & role management",
        "✓ Production deployment with 499+ active users"
    ])

# Slide 17: Closing Slide
add_title_slide(prs,
    "Thank You",
    "Inventory Management System - Enterprise Solution")

# Save presentation
output_path = "IMS_System_Presentation.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation created successfully!")
print(f"📊 File saved as: {output_path}")
print(f"📈 Total slides: {len(prs.slides)}")
