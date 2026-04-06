from pptx import Presentation
from pathlib import Path

ppt_path = Path("Presentation-Inventory Management System (IMS) - 10-03-2026.pptx")
prs = Presentation(ppt_path)

print(f"FILE: {ppt_path.name}")
print(f"SLIDES: {len(prs.slides)}")
print("=" * 80)

for i, slide in enumerate(prs.slides, start=1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            t = " ".join(shape.text.strip().split())
            if t:
                texts.append(t)

    total_chars = sum(len(t) for t in texts)
    total_blocks = len(texts)

    title = "(No clear title)"
    if texts:
        title = texts[0][:140]

    print(f"Slide {i}: blocks={total_blocks}, chars={total_chars}")
    print(f"  Title guess: {title}")

    preview = " | ".join(texts[:4])
    if len(preview) > 320:
        preview = preview[:320] + "..."
    print(f"  Preview: {preview}")
    print("-" * 80)
