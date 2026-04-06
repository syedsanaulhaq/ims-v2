from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 16:9 presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color system
NAVY = RGBColor(18, 42, 76)
BLUE = RGBColor(33, 88, 164)
TEAL = RGBColor(20, 144, 125)
ORANGE = RGBColor(222, 122, 34)
GREEN = RGBColor(46, 140, 87)
LIGHT_BG = RGBColor(245, 248, 252)
MID_BG = RGBColor(232, 239, 248)
TEXT = RGBColor(31, 37, 46)
MUTED = RGBColor(95, 108, 125)
WHITE = RGBColor(255, 255, 255)


def add_background(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.6), Inches(0.45))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(11.5), Inches(0.28))
        stf = st.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = RGBColor(207, 220, 238)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # Visual blocks
    block1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.2), Inches(12.0), Inches(4.8))
    block1.fill.solid()
    block1.fill.fore_color.rgb = RGBColor(24, 53, 92)
    block1.line.fill.background()

    flow = ["Tender In", "Stock Acquisition", "Stock Issuance"]
    x_positions = [1.2, 4.95, 8.75]
    colors = [BLUE, TEAL, ORANGE]

    for i, label in enumerate(flow):
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(3.1), Inches(3.2), Inches(1.1))
        node.fill.solid()
        node.fill.fore_color.rgb = colors[i]
        node.line.fill.background()
        tf = node.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE

        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x_positions[i] + 3.25), Inches(3.35), Inches(0.45), Inches(0.6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = WHITE
            arrow.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.7), Inches(1.1))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    s = slide.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(11.7), Inches(0.55))
    sf = s.text_frame
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(207, 220, 238)
    sp.alignment = PP_ALIGN.CENTER


def add_section_divider(title, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = MID_BG
    stripe.line.fill.background()

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.2), Inches(13.33), Inches(3.1))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.95), Inches(11.8), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    c = slide.shapes.add_textbox(Inches(0.8), Inches(3.85), Inches(11.8), Inches(0.4))
    cf = c.text_frame
    cp = cf.paragraphs[0]
    cp.text = caption
    cp.alignment = PP_ALIGN.CENTER
    cp.font.size = Pt(15)
    cp.font.color.rgb = RGBColor(205, 217, 236)


def add_three_step_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "System Workflow Overview", "Simple business flow used in operations and client reporting")

    step_titles = ["1. Tender In", "2. Stock Acquisition", "3. Stock Issuance"]
    step_sub = [
        "Create tender and finalize vendor path",
        "Receive delivery and post stock in",
        "Approve requests and issue stock out"
    ]
    colors = [BLUE, TEAL, ORANGE]
    x = [0.8, 4.7, 8.6]

    for i in range(3):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x[i]), Inches(2.0), Inches(3.9), Inches(3.7))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = MID_BG

        cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x[i]), Inches(2.0), Inches(3.9), Inches(0.65))
        cap.fill.solid()
        cap.fill.fore_color.rgb = colors[i]
        cap.line.fill.background()

        tt = slide.shapes.add_textbox(Inches(x[i] + 0.15), Inches(2.12), Inches(3.6), Inches(0.4))
        tf = tt.text_frame
        p = tf.paragraphs[0]
        p.text = step_titles[i]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        desc = slide.shapes.add_textbox(Inches(x[i] + 0.2), Inches(2.85), Inches(3.5), Inches(2.5))
        df = desc.text_frame
        dp = df.paragraphs[0]
        dp.text = step_sub[i]
        dp.font.size = Pt(18)
        dp.font.bold = True
        dp.font.color.rgb = TEXT
        dp.alignment = PP_ALIGN.CENTER

        if i < 2:
            ar = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x[i] + 3.95), Inches(3.45), Inches(0.55), Inches(0.75))
            ar.fill.solid()
            ar.fill.fore_color.rgb = NAVY
            ar.line.fill.background()


def add_flow_slide(title, color, steps, emphasis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, title)

    # left visual process lane
    lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.35), Inches(8.35), Inches(5.75))
    lane.fill.solid()
    lane.fill.fore_color.rgb = WHITE
    lane.line.color.rgb = MID_BG

    # right key message
    msg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.35), Inches(3.6), Inches(5.75))
    msg.fill.solid()
    msg.fill.fore_color.rgb = RGBColor(238, 244, 252)
    msg.line.color.rgb = MID_BG

    my = slide.shapes.add_textbox(Inches(9.45), Inches(1.65), Inches(3.15), Inches(0.45))
    myf = my.text_frame
    p = myf.paragraphs[0]
    p.text = "Client Message"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY

    mbody = slide.shapes.add_textbox(Inches(9.45), Inches(2.15), Inches(3.15), Inches(4.6))
    mbf = mbody.text_frame
    bp = mbf.paragraphs[0]
    bp.text = emphasis
    bp.font.size = Pt(16)
    bp.font.bold = True
    bp.font.color.rgb = TEXT

    y = 1.8
    for i, step in enumerate(steps):
        n = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(y), Inches(7.45), Inches(0.9))
        n.fill.solid()
        n.fill.fore_color.rgb = WHITE
        n.line.color.rgb = color

        idx = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05), Inches(y + 0.16), Inches(0.56), Inches(0.56))
        idx.fill.solid()
        idx.fill.fore_color.rgb = color
        idx.line.fill.background()

        it = idx.text_frame
        it.clear()
        ip = it.paragraphs[0]
        ip.text = str(i + 1)
        ip.alignment = PP_ALIGN.CENTER
        ip.font.bold = True
        ip.font.size = Pt(14)
        ip.font.color.rgb = WHITE

        txt = slide.shapes.add_textbox(Inches(1.75), Inches(y + 0.2), Inches(6.5), Inches(0.5))
        tf = txt.text_frame
        tp = tf.paragraphs[0]
        tp.text = step
        tp.font.size = Pt(16)
        tp.font.color.rgb = TEXT

        if i < len(steps) - 1:
            connector = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.2), Inches(y + 0.92), Inches(0.4), Inches(0.28))
            connector.fill.solid()
            connector.fill.fore_color.rgb = color
            connector.line.fill.background()

        y += 1.1


def add_status_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Current Delivery Status", "What is completed with Admin team and what is next")

    # Big progress bar visual
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.0))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(223, 231, 242)
    bar_bg.line.fill.background()

    seg1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(3.8), Inches(1.0))
    seg1.fill.solid()
    seg1.fill.fore_color.rgb = BLUE
    seg1.line.fill.background()

    seg2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.55), Inches(1.9), Inches(3.8), Inches(1.0))
    seg2.fill.solid()
    seg2.fill.fore_color.rgb = TEAL
    seg2.line.fill.background()

    seg3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.9), Inches(4.2), Inches(1.0))
    seg3.fill.solid()
    seg3.fill.fore_color.rgb = RGBColor(238, 240, 244)
    seg3.line.fill.background()

    labels = [
        ("Tender In", 1.9, WHITE),
        ("Stock In", 5.65, WHITE),
        ("Stock Out (Next)", 9.35, TEXT),
    ]

    for text, xpos, c in labels:
        lb = slide.shapes.add_textbox(Inches(xpos), Inches(2.2), Inches(2.4), Inches(0.4))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = c

    # Completed card
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.45), Inches(5.75), Inches(3.0))
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = MID_BG

    head1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.45), Inches(5.75), Inches(0.55))
    head1.fill.solid()
    head1.fill.fore_color.rgb = GREEN
    head1.line.fill.background()

    h1 = slide.shapes.add_textbox(Inches(1.05), Inches(3.58), Inches(5.2), Inches(0.3))
    h1p = h1.text_frame.paragraphs[0]
    h1p.text = "Completed with Admin Team"
    h1p.font.size = Pt(14)
    h1p.font.bold = True
    h1p.font.color.rgb = WHITE

    body1 = slide.shapes.add_textbox(Inches(1.05), Inches(4.15), Inches(5.3), Inches(2.2))
    btf = body1.text_frame
    points = [
        "Tender entries created and validated",
        "Delivery receiving tested",
        "Stock acquisition records generated",
        "Inventory updates confirmed"
    ]
    for i, t in enumerate(points):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = "- " + t
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT

    # Next card
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(3.45), Inches(5.75), Inches(3.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = MID_BG

    head2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.75), Inches(3.45), Inches(5.75), Inches(0.55))
    head2.fill.solid()
    head2.fill.fore_color.rgb = ORANGE
    head2.line.fill.background()

    h2 = slide.shapes.add_textbox(Inches(7.0), Inches(3.58), Inches(5.2), Inches(0.3))
    h2p = h2.text_frame.paragraphs[0]
    h2p.text = "In Progress / Next Step"
    h2p.font.size = Pt(14)
    h2p.font.bold = True
    h2p.font.color.rgb = WHITE

    body2 = slide.shapes.add_textbox(Inches(7.0), Inches(4.15), Inches(5.3), Inches(2.2))
    ctf = body2.text_frame
    points2 = [
        "Stock issuance (stock out) execution",
        "Approval and dispatch sequence",
        "Final operational handover",
        "Client acceptance walkthrough"
    ]
    for i, t in enumerate(points2):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = "- " + t
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT


def add_closing_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.7), Inches(1.0))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "IMS Workflow Demonstration"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(11.1), Inches(1.0))
    sf = sub.text_frame
    sp = sf.paragraphs[0]
    sp.text = "Tender In  ->  Stock In  ->  Stock Out"
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = Pt(24)
    sp.font.bold = True
    sp.font.color.rgb = RGBColor(206, 222, 245)

    note = slide.shapes.add_textbox(Inches(2.0), Inches(4.3), Inches(9.3), Inches(1.4))
    nf = note.text_frame
    np = nf.paragraphs[0]
    np.text = "Tender to Stock In is validated with Admin test entries.\nStock Out is the active next delivery step."
    np.alignment = PP_ALIGN.CENTER
    np.font.size = Pt(16)
    np.font.color.rgb = WHITE


# Build deck
add_title_slide(
    "Inventory Management System",
    "Client Workflow Presentation - Visual Overview"
)
add_three_step_overview()
add_section_divider("Step 1: Tender In", "Where every inventory cycle starts")
add_flow_slide(
    "Tender Creation and Finalization",
    BLUE,
    [
        "Create tender with required items and quantities",
        "Engage vendors and evaluate responses",
        "Finalize tender and prepare PO path"
    ],
    "This stage ensures demand is formally approved before stock movement starts."
)
add_section_divider("Step 2: Stock Acquisition (Stock In)", "Convert approved procurement into available inventory")
add_flow_slide(
    "Delivery Receiving and Stock Posting",
    TEAL,
    [
        "Receive delivery against purchase order",
        "Validate quantity and quality during receiving",
        "Create stock acquisition and update inventory"
    ],
    "This stage is completed and tested end-to-end with Admin team entries."
)
add_section_divider("Step 3: Stock Issuance (Stock Out)", "Release inventory to requesting users or wings")
add_flow_slide(
    "Issuance Flow (Next Active Step)",
    ORANGE,
    [
        "User raises stock request",
        "Approval flow confirms request eligibility",
        "Approved quantity is issued and deducted"
    ],
    "Stock out is the next execution step currently in progress."
)
add_status_slide()
add_closing_slide()

output_path = "IMS-Client-Workflow-Visual-Presentation.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
print(f"Total slides: {len(prs.slides)}")
